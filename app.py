"""
리더십 진단 보고서 자동화 툴 v3
수정사항:
  1. 엑셀 - 6대 역량 / 8대 기술 평균값을 Python에서 직접 계산해 셀에 주입 (수식 의존 제거)
  2. PPT  - 슬라이드 복제 시 차트 관계(rId) 포함 완전 복제 방식으로 교체
           (참고 코드의 pptx 내부 XML + _part 복사 패턴 적용)
"""

import streamlit as st
import pandas as pd
import openpyxl
from openpyxl import load_workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
import copy
import io
import zipfile
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.opc.packuri import PackURI
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

# ══════════════════════════════════════════════════════════════════
# 1. 매핑 정의
# ══════════════════════════════════════════════════════════════════

COMPETENCY_MAP = {
    "Position":     [6, 7, 14, 15, 19, 28],
    "Personality":  [2, 10, 11, 18, 21, 25],
    "Relationship": [3, 4, 20, 22, 27, 29],
    "Results":      [5, 13, 26, 30],
    "Development":  [1, 9, 17, 24],
    "Principles":   [8, 12, 16, 23],
}

SKILL_MAP = {
    "우호성":     [1, 9, 17, 24],
    "동기유발":   [2, 10, 18, 25],
    "자문":       [3, 11],
    "협력제휴":   [4, 12, 19, 26],
    "협상거래":   [5, 13, 20, 27],
    "합리적설득": [6, 14, 21, 28],
    "합법화":     [7, 15, 22, 29],
    "강요":       [8, 16, 23, 30],
}

SOFT_SKILLS = ["우호성", "동기유발", "자문"]
HARD_SKILLS = ["협력제휴", "협상거래", "합리적설득", "합법화", "강요"]

# ══════════════════════════════════════════════════════════════════
# 2. 점수 계산
# ══════════════════════════════════════════════════════════════════

def calc_avg(scores: dict, q_list: list) -> float:
    vals = [scores[q] for q in q_list if q in scores]
    return round(sum(vals) / len(vals), 2) if vals else 0.0


def compute_person(scores: dict) -> dict:
    competency = {k: calc_avg(scores, v) for k, v in COMPETENCY_MAP.items()}
    skill_raw  = {k: calc_avg(scores, v) for k, v in SKILL_MAP.items()}
    soft_avg   = round(sum(skill_raw[k] for k in SOFT_SKILLS) / len(SOFT_SKILLS), 2)
    hard_avg   = round(sum(skill_raw[k] for k in HARD_SKILLS) / len(HARD_SKILLS), 2)
    return {"competency": competency, "skill_raw": skill_raw,
            "soft_avg": soft_avg, "hard_avg": hard_avg}

# ══════════════════════════════════════════════════════════════════
# 3. 입력 파싱
# ══════════════════════════════════════════════════════════════════

def parse_response_excel(file) -> list:
    """
    A열(idx 0): 타임스탬프, B열(idx 1): 성함, C열(idx 2)~AF열: Q1~Q30
    col_idx(0-based) = Q + 1  (Q1 → iloc[2] = C열)
    """
    df = pd.read_excel(file, header=0)
    people = []
    for _, row in df.iterrows():
        name = str(row.iloc[1]).strip()
        if not name or name.lower() in ("nan", ""):
            continue
        scores = {}
        for q in range(1, 31):
            try:
                scores[q] = float(row.iloc[q + 1])
            except Exception:
                scores[q] = 0.0
        people.append({"name": name, "scores": scores})
    return people

# ══════════════════════════════════════════════════════════════════
# 4-A. 엑셀 템플릿 자동 생성
# ══════════════════════════════════════════════════════════════════

def make_default_excel_template() -> io.BytesIO:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Template"

    DARK_BLUE  = "1F3864"
    MID_BLUE   = "2E75B6"
    LIGHT_BLUE = "BDD7EE"
    ACCENT     = "ED7D31"
    WHITE      = "FFFFFF"
    GRAY_BG    = "F2F2F2"
    YELLOW_BG  = "FFF2CC"

    thin   = Side(border_style="thin", color="BFBFBF")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    def hdr(coord, val, bg=DARK_BLUE, fg=WHITE, bold=True, size=11, align="center"):
        c = ws[coord]
        c.value = val
        c.font = Font(bold=bold, color=fg, size=size, name="맑은 고딕")
        c.fill = PatternFill("solid", fgColor=bg)
        c.alignment = Alignment(horizontal=align, vertical="center", wrap_text=True)
        c.border = border

    def dat(coord, val, bg=WHITE, bold=False, align="center"):
        c = ws[coord]
        c.value = val
        c.font = Font(bold=bold, size=10, name="맑은 고딕")
        c.fill = PatternFill("solid", fgColor=bg)
        c.alignment = Alignment(horizontal=align, vertical="center")
        c.border = border

    def num_cell(coord, val, bg=WHITE, bold=False):
        c = ws[coord]
        c.value = val
        c.font = Font(bold=bold, size=10, name="맑은 고딕")
        c.fill = PatternFill("solid", fgColor=bg)
        c.alignment = Alignment(horizontal="center", vertical="center")
        c.number_format = "0.00"
        c.border = border

    for col, w in {"A":13,"B":7,"C":10,"D":2,"E":18,"F":12,"G":2,"H":18,"I":12}.items():
        ws.column_dimensions[col].width = w
    ws.row_dimensions[1].height = 30
    ws.row_dimensions[2].height = 8
    ws.row_dimensions[3].height = 22
    for r in range(4, 37):
        ws.row_dimensions[r].height = 18

    # 타이틀
    hdr("A1", "리더십 진단 보고서", bg=DARK_BLUE, size=14)
    hdr("B1", "응답자", bg=MID_BLUE)
    c = ws["C1"]
    c.value = ""
    c.font = Font(bold=True, size=12, color=DARK_BLUE, name="맑은 고딕")
    c.fill = PatternFill("solid", fgColor=LIGHT_BLUE)
    c.alignment = Alignment(horizontal="center", vertical="center")
    c.border = border

    # 문항 헤더
    hdr("A3", "문항번호", bg=MID_BLUE)
    hdr("B3", "Q#",      bg=MID_BLUE)
    hdr("C3", "점수",    bg=MID_BLUE)

    # Q1~Q30 레이블 + 점수 셀
    for q in range(1, 31):
        row = q + 3
        bg = GRAY_BG if q % 2 == 0 else WHITE
        dat(f"A{row}", f"문항 {q}", bg=bg)
        dat(f"B{row}", q,           bg=bg)
        sc = ws[f"C{row}"]
        sc.value = None
        sc.font = Font(size=10, name="맑은 고딕")
        sc.fill = PatternFill("solid", fgColor=YELLOW_BG)
        sc.alignment = Alignment(horizontal="center", vertical="center")
        sc.border = border

    # 6대 역량 헤더 + 빈 평균칸 (값은 build_excel에서 주입)
    hdr("E3", "6대 역량",  bg=DARK_BLUE)
    hdr("F3", "평균점수",  bg=DARK_BLUE)
    for i, key in enumerate(COMPETENCY_MAP):
        row = i + 4
        bg = GRAY_BG if i % 2 == 0 else WHITE
        dat(f"E{row}", key, bg=bg, align="left")
        num_cell(f"F{row}", None, bg=bg)

    # 8대 기술 헤더 + 빈 평균칸
    hdr("H3", "8대 기술",  bg=ACCENT)
    hdr("I3", "평균점수",  bg=ACCENT)
    skill_keys = list(SKILL_MAP.keys())
    for i, key in enumerate(skill_keys):
        row = i + 4
        is_soft = key in SOFT_SKILLS
        bg = "E2EFDA" if is_soft else (GRAY_BG if i % 2 == 0 else WHITE)
        dat(f"H{row}", key, bg=bg, align="left")
        num_cell(f"I{row}", None, bg=bg)

    sr = len(skill_keys) + 4
    hdr(f"H{sr}", "소프트스킬 평균", bg="70AD47", size=10)
    num_cell(f"I{sr}", None, bg="A9D18E", bold=True)

    hr_ = sr + 1
    hdr(f"H{hr_}", "하드스킬 평균", bg=ACCENT, size=10)
    num_cell(f"I{hr_}", None, bg="F4B183", bold=True)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf

# ══════════════════════════════════════════════════════════════════
# 4-B. 엑셀 출력 생성  ★ 평균값 직접 계산해 주입
# ══════════════════════════════════════════════════════════════════

def _copy_sheet(wb_dest, src_ws, new_title: str):
    """워크시트를 다른 워크북으로 안전하게 복사. ws 반환."""
    ws = wb_dest.create_sheet(title=new_title)
    for col_letter, cd in src_ws.column_dimensions.items():
        ws.column_dimensions[col_letter].width = cd.width
    for row_num, rd in src_ws.row_dimensions.items():
        ws.row_dimensions[row_num].height = rd.height
    for row in src_ws.iter_rows():
        for cell in row:
            nc = ws.cell(row=cell.row, column=cell.column)
            nc.value = cell.value
            if cell.has_style:
                nc.font          = copy.copy(cell.font)
                nc.border        = copy.copy(cell.border)
                nc.fill          = copy.copy(cell.fill)
                nc.number_format = cell.number_format
                nc.protection    = copy.copy(cell.protection)
                nc.alignment     = copy.copy(cell.alignment)
    for merge in src_ws.merged_cells.ranges:
        ws.merge_cells(str(merge))
    return ws   # ← 반드시 return


def build_excel(people: list, template_src) -> bytes:
    """
    각 응답자 시트 생성:
      C1      = 성함
      C4:C33  = Q1~Q30 점수 (숫자)
      F4:F9   = 6대 역량 평균 (Python 계산값 직접 주입)
      I4:I11  = 8대 기술 평균 (Python 계산값 직접 주입)
      I12     = 소프트스킬 평균
      I13     = 하드스킬 평균
    """
    if hasattr(template_src, "read"):
        raw = template_src.read()
        template_src = io.BytesIO(raw)

    wb_out = openpyxl.Workbook()
    wb_out.remove(wb_out.active)

    skill_keys = list(SKILL_MAP.keys())

    for person in people:
        template_src.seek(0)
        wb_tpl = load_workbook(template_src)
        src_ws = wb_tpl.worksheets[0]

        safe_name = person["name"][:31]
        ws = _copy_sheet(wb_out, src_ws, safe_name)

        result = compute_person(person["scores"])

        # ① 성함 (C1)
        ws.cell(row=1, column=3).value = person["name"]

        # ② Q1~Q30 점수 (C4:C33)
        for q in range(1, 31):
            ws.cell(row=q + 3, column=3).value = person["scores"].get(q, 0)

        # ③ 6대 역량 평균 (F4:F9) — 계산값 직접 주입
        for i, (key, avg) in enumerate(result["competency"].items()):
            c = ws.cell(row=i + 4, column=6)
            c.value = avg
            c.number_format = "0.00"

        # ④ 8대 기술 평균 (I4:I11) — 계산값 직접 주입
        for i, key in enumerate(skill_keys):
            c = ws.cell(row=i + 4, column=9)
            c.value = result["skill_raw"][key]
            c.number_format = "0.00"

        # ⑤ 소프트스킬 / 하드스킬 평균 (I12, I13)
        sr_row = len(skill_keys) + 4
        c_soft = ws.cell(row=sr_row,     column=9)
        c_hard = ws.cell(row=sr_row + 1, column=9)
        c_soft.value = result["soft_avg"]; c_soft.number_format = "0.00"
        c_hard.value = result["hard_avg"]; c_hard.number_format = "0.00"

    buf = io.BytesIO()
    wb_out.save(buf)
    buf.seek(0)
    return buf.read()

# ══════════════════════════════════════════════════════════════════
# 5-A. PPT 템플릿 자동 생성
# ══════════════════════════════════════════════════════════════════

def make_default_ppt_template() -> io.BytesIO:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)

    def add_txt(l, t, w, h, text, size=12, bold=False,
                color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT, name=None):
        box = slide.shapes.add_textbox(l, t, w, h)
        if name: box.name = name
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return box

    # 상단 타이틀 바
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x1F, 0x38, 0x64)
    bar.line.fill.background()
    tf = bar.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].add_run()
    run.text = "  리더십 진단 보고서"
    run.font.size = Pt(22); run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    add_txt(Inches(9.5), Inches(0.12), Inches(3.5), Inches(0.65),
            "{{NAME}}", size=18, bold=True,
            color=RGBColor(0xFF, 0xD7, 0x00),
            align=PP_ALIGN.RIGHT, name="name_label")

    add_txt(Inches(0.2), Inches(1.0), Inches(4.4), Inches(0.4),
            "6대 역량 (Phase)", size=13, bold=True, color=RGBColor(0x1F,0x38,0x64))
    add_txt(Inches(4.8), Inches(1.0), Inches(4.4), Inches(0.4),
            "8대 영향력 기술 (Strategy)", size=13, bold=True, color=RGBColor(0xED,0x7D,0x31))
    add_txt(Inches(9.5), Inches(1.0), Inches(3.6), Inches(0.4),
            "역량 레이더 차트", size=13, bold=True, color=RGBColor(0x1F,0x38,0x64))

    # table_phase: 헤더1 + 데이터6 = 7행
    tbl_p = slide.shapes.add_table(7, 2, Inches(0.2), Inches(1.45), Inches(4.4), Inches(5.8))
    tbl_p.name = "table_phase"
    t = tbl_p.table
    for ci, h in enumerate(["항목명", "평균점수"]):
        cell = t.cell(0, ci)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1F,0x38,0x64)
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]; run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); run.font.size = Pt(11)
    for ri, key in enumerate(COMPETENCY_MAP.keys(), start=1):
        c0 = t.cell(ri, 0); c0.text = key
        c1 = t.cell(ri, 1); c1.text = ""
        bg = RGBColor(0xBD,0xD7,0xEE) if ri % 2 == 0 else RGBColor(0xFF,0xFF,0xFF)
        for ci2, c in enumerate([c0, c1]):
            c.fill.solid(); c.fill.fore_color.rgb = bg
            pp = c.text_frame.paragraphs[0]
            pp.alignment = PP_ALIGN.CENTER if ci2 == 1 else PP_ALIGN.LEFT
            if pp.runs: pp.runs[0].font.size = Pt(10)

    # table_strategy: 헤더1 + 소프트3 + 소프트평균1 + 하드5 + 하드평균1 = 11행
    tbl_s = slide.shapes.add_table(11, 2, Inches(4.8), Inches(1.45), Inches(4.4), Inches(5.8))
    tbl_s.name = "table_strategy"
    t2 = tbl_s.table
    for ci, h in enumerate(["항목명", "평균점수"]):
        cell = t2.cell(0, ci)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xED,0x7D,0x31)
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]; run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); run.font.size = Pt(11)
    strat_rows = SOFT_SKILLS + ["소프트스킬 평균"] + HARD_SKILLS + ["하드스킬 평균"]
    for ri, label in enumerate(strat_rows, start=1):
        is_avg  = "평균" in label
        is_soft = label in SOFT_SKILLS or label == "소프트스킬 평균"
        c0 = t2.cell(ri, 0); c0.text = label
        c1 = t2.cell(ri, 1); c1.text = ""
        if is_avg:
            bg = RGBColor(0xA9,0xD1,0x8E) if is_soft else RGBColor(0xF4,0xB1,0x83)
        elif is_soft:
            bg = RGBColor(0xE2,0xEF,0xDA)
        else:
            bg = RGBColor(0xF2,0xF2,0xF2) if ri % 2 == 0 else RGBColor(0xFF,0xFF,0xFF)
        for ci2, c in enumerate([c0, c1]):
            c.fill.solid(); c.fill.fore_color.rgb = bg
            pp = c.text_frame.paragraphs[0]
            pp.alignment = PP_ALIGN.CENTER if ci2 == 1 else PP_ALIGN.LEFT
            if pp.runs:
                pp.runs[0].font.size = Pt(10)
                if is_avg: pp.runs[0].font.bold = True

    # chart_phase: 레이더 차트
    cd = ChartData()
    cd.categories = list(COMPETENCY_MAP.keys())
    cd.add_series("역량 점수", [3.0] * 6)
    cs = slide.shapes.add_chart(
        XL_CHART_TYPE.RADAR_FILLED,
        Inches(9.5), Inches(1.45), Inches(3.6), Inches(5.8), cd
    )
    cs.name = "chart_phase"

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

# ══════════════════════════════════════════════════════════════════
# 5-B. PPT 슬라이드 완전 복제  ★ 차트 관계(rId) 포함 복사
# ══════════════════════════════════════════════════════════════════

def _clone_slide_full(prs: Presentation, src_slide_index: int = 0):
    """
    참고 코드 패턴: 슬라이드 XML + 모든 관계(차트, 이미지 등) 포함 완전 복제.
    python-pptx 내부 Part 구조를 직접 복사하여 rId 누락 문제 해결.
    """
    src_slide = prs.slides[src_slide_index]

    # 새 슬라이드 추가 (빈 레이아웃)
    layout    = src_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # 새 슬라이드의 XML을 원본 XML로 완전 교체
    new_slide._element = copy.deepcopy(src_slide._element)

    # ── 관계(relationships) 복사 ──
    # 원본 슬라이드의 모든 관계를 새 슬라이드에 동일하게 복사
    src_part = src_slide.part
    new_part = new_slide.part

    for rel in src_part.rels.values():
        # 이미 존재하는 관계는 건너뜀
        if rel.reltype in new_part.rels:
            continue
        if rel.is_external:
            new_part.rels._rels[rel.rId] = rel
        else:
            # 내부 파트(차트, 이미지 등) deep copy
            target_part = rel._target
            new_part.rels._rels[rel.rId] = rel.__class__(
                rel.rId, rel.reltype, target_part, rel.is_external
            )

    return new_slide


def _replace_text_in_shape(shape, old: str, new_val: str):
    """shape 내 텍스트에서 old → new_val 치환."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new_val)


def _set_table_data(shape, rows_data: list):
    """표 1행~에 데이터 주입 (0행 헤더 유지)."""
    tbl = shape.table
    for r_idx, (label, val) in enumerate(rows_data, start=1):
        if r_idx >= len(tbl.rows):
            break
        tbl.cell(r_idx, 0).text = str(label)
        tbl.cell(r_idx, 1).text = f"{val:.2f}" if isinstance(val, float) else str(val)


def _update_chart(shape, competency: dict):
    try:
        cd = ChartData()
        cd.categories = list(competency.keys())
        cd.add_series("역량 점수", list(competency.values()))
        shape.chart.replace_data(cd)
    except Exception:
        pass


def _fill_slide(slide, name: str, result: dict):
    competency = result["competency"]
    skill_raw  = result["skill_raw"]
    soft_avg   = result["soft_avg"]
    hard_avg   = result["hard_avg"]

    for shape in slide.shapes:
        sname = shape.name.lower()

        _replace_text_in_shape(shape, "{{NAME}}", name)
        _replace_text_in_shape(shape, "{{name}}", name)

        if "table_phase" in sname and shape.has_table:
            _set_table_data(shape, list(competency.items()))

        elif "table_strategy" in sname and shape.has_table:
            rows = [(k, skill_raw[k]) for k in SOFT_SKILLS]
            rows.append(("소프트스킬 평균", soft_avg))
            rows += [(k, skill_raw[k]) for k in HARD_SKILLS]
            rows.append(("하드스킬 평균", hard_avg))
            _set_table_data(shape, rows)

        elif "chart_phase" in sname and shape.has_chart:
            _update_chart(shape, competency)


def build_ppt(people: list, template_src) -> bytes:
    """
    핵심 전략:
      1. 첫 번째 사람 → 템플릿 슬라이드[0]에 직접 주입
      2. 두 번째 이후 → _clone_slide_full()로 슬라이드[0] 완전 복제 후 주입
         (복제는 항상 원본 슬라이드[0] 기준으로 반복)
    """
    if hasattr(template_src, "read"):
        raw = template_src.read()
        template_src = io.BytesIO(raw)

    template_src.seek(0)
    prs = Presentation(template_src)

    for i, person in enumerate(people):
        result = compute_person(person["scores"])
        if i == 0:
            slide = prs.slides[0]
        else:
            # 항상 슬라이드[0] (원본)을 기준으로 복제
            slide = _clone_slide_full(prs, src_slide_index=0)
        _fill_slide(slide, person["name"], result)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ══════════════════════════════════════════════════════════════════
# 6. Streamlit UI
# ══════════════════════════════════════════════════════════════════

st.set_page_config(page_title="리더십 진단 보고서 자동화", layout="wide")

st.markdown("""
<style>
.main-title{font-size:2rem;font-weight:800;color:#1F3864;}
.sub{font-size:1rem;color:#666;margin-bottom:.5rem;}
.sec{font-size:1.05rem;font-weight:700;color:#2E75B6;margin-top:1.2rem;margin-bottom:.3rem;}
</style>""", unsafe_allow_html=True)

st.markdown('<div class="main-title">📊 리더십 진단 보고서 자동화 툴</div>', unsafe_allow_html=True)
st.markdown('<div class="sub">구글 폼 응답 엑셀 → 개인별 엑셀 (평균값 포함) + 응답자별 슬라이드 PPT 자동 생성</div>',
            unsafe_allow_html=True)
st.markdown("---")

st.markdown('<div class="sec">① 구글 폼 응답 엑셀 (필수)</div>', unsafe_allow_html=True)
response_file = st.file_uploader(
    "A열: 타임스탬프 / B열: 성함 / C~AF열: Q1~Q30 점수",
    type=["xlsx", "xls"], key="response"
)

st.markdown('<div class="sec">② 엑셀 템플릿 (선택 – 미업로드 시 자동 생성)</div>', unsafe_allow_html=True)
excel_tpl_file = st.file_uploader(
    "C1=성함, C4:C33=점수, F열=6대역량 평균, I열=8대기술 평균 위치인 양식",
    type=["xlsx"], key="excel_tpl"
)

st.markdown('<div class="sec">③ PPT 템플릿 (선택 – 미업로드 시 자동 생성)</div>', unsafe_allow_html=True)
ppt_tpl_file = st.file_uploader(
    "{{NAME}} · table_phase · table_strategy · chart_phase 개체 포함 1슬라이드",
    type=["pptx"], key="ppt_tpl"
)

with st.expander("📥 기본 템플릿 미리 다운로드 (편집 후 업로드 가능)"):
    ca, cb = st.columns(2)
    with ca:
        st.download_button("⬇️ 기본 엑셀 템플릿",
            data=make_default_excel_template(), file_name="excel_template.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True)
    with cb:
        st.download_button("⬇️ 기본 PPT 템플릿",
            data=make_default_ppt_template(), file_name="template.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True)

st.markdown("---")

if st.button("🚀 보고서 생성", type="primary", use_container_width=True):
    if not response_file:
        st.error("❌ 구글 폼 응답 엑셀을 업로드해주세요.")
        st.stop()

    with st.spinner("📂 응답 데이터 파싱 중..."):
        try:
            people = parse_response_excel(response_file)
        except Exception as e:
            st.error(f"❌ 응답 엑셀 파싱 실패: {e}")
            st.stop()

    if not people:
        st.error("❌ 응답자 데이터가 비어있습니다.")
        st.stop()

    st.success(f"✅ {len(people)}명 응답 데이터 파싱 완료")

    with st.expander(f"📋 응답자 미리보기 ({len(people)}명)"):
        preview = []
        for p in people:
            r = compute_person(p["scores"])
            row = {"성함": p["name"]}
            row.update({k: f"{v:.2f}" for k, v in r["competency"].items()})
            row["소프트스킬 평균"] = f"{r['soft_avg']:.2f}"
            row["하드스킬 평균"]   = f"{r['hard_avg']:.2f}"
            preview.append(row)
        st.dataframe(pd.DataFrame(preview), use_container_width=True)

    excel_src = excel_tpl_file if excel_tpl_file else make_default_excel_template()
    ppt_src   = ppt_tpl_file   if ppt_tpl_file   else make_default_ppt_template()
    if not excel_tpl_file: st.info("ℹ️ 엑셀 템플릿 미업로드 → 기본 템플릿 자동 사용")
    if not ppt_tpl_file:   st.info("ℹ️ PPT 템플릿 미업로드 → 기본 템플릿 자동 사용")

    with st.spinner("📊 개인별 엑셀 생성 중 (평균값 포함)..."):
        try:
            excel_bytes = build_excel(people, excel_src)
        except Exception as e:
            st.error(f"❌ 엑셀 생성 실패: {e}")
            st.exception(e); st.stop()

    with st.spinner(f"📑 PPT 생성 중 (총 {len(people)}슬라이드)..."):
        try:
            ppt_bytes = build_ppt(people, ppt_src)
        except Exception as e:
            st.error(f"❌ PPT 생성 실패: {e}")
            st.exception(e); st.stop()

    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("리더십진단_개인별.xlsx", excel_bytes)
        zf.writestr("리더십진단_통합.pptx",  ppt_bytes)
    zip_buf.seek(0)

    st.balloons()
    st.success(f"🎉 완료! 엑셀 {len(people)}시트 + PPT {len(people)}슬라이드 생성")

    d1, d2, d3 = st.columns(3)
    with d1:
        st.download_button("⬇️ ZIP (전체)", data=zip_buf,
            file_name="리더십진단_결과.zip", mime="application/zip",
            use_container_width=True)
    with d2:
        st.download_button("⬇️ 엑셀 (개인별)", data=excel_bytes,
            file_name="리더십진단_개인별.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True)
    with d3:
        st.download_button("⬇️ PPT (통합)", data=ppt_bytes,
            file_name="리더십진단_통합.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True)
